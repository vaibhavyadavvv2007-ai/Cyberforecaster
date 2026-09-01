"""Build the SIH26153 idea presentation from the official SIH template.

Fills SIH2026-IDEA-Presentation-Format.pptx (7 slides) with our content and
removes the instructions slide, producing SIH26153_Idea_Presentation.pptx
(6 slides, per the template's own limit). The template's section pointers
are kept verbatim as bold headings, as required ("without changing the idea
details pointers").

Every number quoted below comes from a repo artifact (models/*.json, the
rehearsal logs in docs/DEMO_RUNBOOK.md §7) — do not edit by feel.

Team ID / Team Name are left as <fill> placeholders — the portal assigns
them; fill before export.
"""
from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "SIH2026-IDEA-Presentation-Format.pptx"
DST = ROOT / "SIH26153_Idea_Presentation.pptx"

HEAD = RGBColor(0x1F, 0x49, 0x7D)   # theme dk2 — the template's heading color
BODY = RGBColor(0x26, 0x26, 0x26)
FONT = "Arial"

H = "heading"
B = "body"


def fill(shape, items, top=None, height=None):
    """Rewrite a shape's text frame: (kind, text) pairs — headings bold in theme blue."""
    if top is not None:
        shape.top = top
    if height is not None:
        shape.height = height
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, (kind, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = FONT
        if kind == H:
            f.size = Pt(17)
            f.bold = True
            f.color.rgb = HEAD
            p.space_before = Pt(12 if i else 0)
            p.space_after = Pt(3)
        else:
            f.size = Pt(13.5)
            f.bold = False
            f.color.rgb = BODY
            p.space_before = Pt(2)
            p.space_after = Pt(2)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


shutil.copy(SRC, DST)
prs = Presentation(str(DST))
S = prs.slides

# ------------------------------------------------------------ slide 1: title
tb9 = next(sh for sh in S[0].shapes if sh.name == "TextBox 9")
lines = [
    "Problem Statement ID: SIH26153",
    "Problem Statement Title: AI-based Network Attack Forecasting from Network Traffic Data",
    "Theme: Blockchain & Cybersecurity  (Organization: NTRO)",
    "PS Category: Software",
    "Team ID: <fill from portal>",
    "Team Name: <registered name>",
]
paras = [p for p in tb9.text_frame.paragraphs if p.runs]
for p, line in zip(paras, lines):
    p.runs[0].text = line
    for extra in p.runs[1:]:
        extra.text = ""
notes(S[0], (
    "Title slide. One-liner if asked: 'We forecast network attacks before they "
    "complete — and show you why.' Fill Team ID / Team Name from the portal "
    "before exporting to PDF."))

# ------------------------------------------------------------ slide 2: idea
title = next(sh for sh in S[1].shapes if sh.name == "Title 1")
ttf = title.text_frame
ttf.clear()
tp = ttf.paragraphs[0]
trun = tp.add_run()
trun.text = "CyberForecaster — forecasting the attack, not just flagging it"
trun.font.name = "Times New Roman"
trun.font.size = Pt(32)
trun.font.bold = True
trun.font.color.rgb = HEAD
tb = next(sh for sh in S[1].shapes if sh.name == "TextBox 8")
fill(tb, [
    (H, "Proposed Solution"),
    (B, "•  A temporal attack forecaster: reads 10 x 30-second network-state windows (18 flow features) and forecasts attack probability for each of the next 5 windows (2.5 min ahead), plus the coming MITRE ATT&CK stage"),
    (B, "•  Working prototype today: FastAPI + Next.js console, fully offline on one laptop — historical attack scenarios AND a live packet sensor (Npcap/scapy) feed the same trained model"),
    (B, "•  Every warning is explained: per-feature attribution (Integrated Gradients) + transparent ATT&CK stage rules"),
    (H, "How it addresses the problem"),
    (B, "•  Implements the PS's world-model requirement directly — P(S t+1 | S t) learned from telemetry and rolled forward K=5 steps; not a static per-flow classifier"),
    (B, "•  Beats the required logistic baseline ~2x on PR-AUC (0.656 vs 0.333) on identical features and split"),
    (B, "•  Explainability (PS asks for SHAP/attention): every prediction ships with feature attribution and an independent rule-based stage cross-check"),
    (H, "Innovation and uniqueness"),
    (B, "•  Forecast, don't just detect: probability trajectory + kill-chain stage trajectory = analyst decision support while the attack is still unfolding"),
    (B, "•  Two engines by design: transparent rules catch instant signatures (recon scan flagged in one window); the LSTM forecasts sustained progression — both verified live over real Wi-Fi"),
    (B, "•  57K parameters · 0.23 MB · 2.6 ms CPU inference → edge-deployable, offline, flow-metadata only (no payloads — privacy-preserving)"),
], top=Inches(1.35), height=Inches(5.6))
notes(S[1], (
    "The demo answers 'does it work' before the jury can ask: we will SHOW this "
    "live. Numbers: threshold 0.561 (validation 5% FPR budget), precision 0.88, "
    "FPR 0.006 on an unseen attack family. IBM 2025 stats on the impact slide."))

# ------------------------------------------------- slide 3: technical approach
tb = next(sh for sh in S[2].shapes if sh.name == "TextBox 8")
fill(tb, [
    (H, "Technologies to be used"),
    (B, "•  Python · PyTorch (2-layer LSTM, multi-task heads) · scapy/Npcap live capture · FastAPI backend · Next.js + TypeScript console · CSE-CIC-IDS2018 benchmark (~6,200 windows, 1,486 attack)"),
    (H, "Methodology and process for implementation"),
    (B, "•  Flow records → 30-s windows x 18 features → chronological 70/15/15 split with boundary purge (no leakage; test split = an attack family absent from training)"),
    (B, "•  log1p + standardize fitted on train only — one shared transform for baseline, LSTM and the live sensor"),
    (B, "•  2-layer LSTM (hidden 64): 5 progression logits + 6-stage head; alert threshold from a 5% false-positive budget on validation"),
    (B, "•  Independent rule engine (6 ATT&CK rules) validated against labelled windows; Integrated Gradients attribution on every prediction"),
], top=Inches(1.30), height=Inches(3.20))
arch = S[2].shapes.add_picture(str(ROOT / "docs/assets/architecture.png"),
                               Inches(0.35), Inches(4.55), width=Inches(7.2))
bench = S[2].shapes.add_picture(str(ROOT / "docs/assets/benchmark.png"),
                                Inches(7.90), Inches(4.95), width=Inches(4.3))
notes(S[2], (
    "Architecture: two inputs (offline dataset + live packets) → one window "
    "builder → two engines (orange rules = instant signatures; blue LSTM = "
    "5-step forecast) → one explained warning. Chart: PR-AUC on the test split "
    "which is an UNSEEN attack family — measures transfer, not memorization. "
    "Same features/split for both models, so the gap isolates temporal "
    "modeling. The live demo runs this exact pipeline on real packets."))

# ------------------------------------------------- slide 4: feasibility
tb = next(sh for sh in S[3].shapes if sh.name == "TextBox 8")
fill(tb, [
    (H, "Analysis of the feasibility of the idea"),
    (B, "•  Already built and verified end-to-end: offline benchmark + live sensor demo; in rehearsal a real attack over Wi-Fi crossed the alert threshold during a sustained UDP sweep while benign traffic stayed LOW"),
    (B, "•  0.23 MB model, 2.6 ms CPU inference, fully offline — runs on a commodity laptop, no cloud dependency"),
    (H, "Potential challenges and risks"),
    (B, "•  Benchmark ↔ production traffic shift (benchmark flows are long-lived aggregates; production sees short transactions)"),
    (B, "•  Class imbalance and rare attack families (~24% of windows contain attack)"),
    (B, "•  False alarms erode analyst trust; no system can warn before the first attack packet exists"),
    (H, "Strategies for overcoming these challenges"),
    (B, "•  Live input conditioned to the model's validated training domain + rule engine on raw features; NetFlow/IPFIX named as the production input"),
    (B, "•  PR-AUC focus, chronological split, high-precision operating point (FPR 0.6%); cross-family test split measures transfer honestly"),
    (B, "•  Every alert explained (attribution + rules); positioned as decision support, not auto-blocking; limitations stated up front"),
    (B, "•  Scale path: windowing at each sensor/tap, one shared 0.23 MB model centrally or at the edge"),
], top=Inches(1.30), height=Inches(5.7))
notes(S[3], (
    "Feasibility is proven, not claimed — the prototype exists and was "
    "rehearsed over real Wi-Fi on Aug 30 (attack forecast climbed 0.03 → 0.905 "
    "HIGH → 0.988 across four sustained windows; benign worst 0.014). Honesty "
    "is the strategy: we state the no-pre-onset-warning limit ourselves."))

# ---------------------------------------------------- slide 5: impact
tb = next(sh for sh in S[4].shapes if sh.name == "TextBox 8")
fill(tb, [
    (H, "Potential impact on the target audience"),
    (B, "•  SOC analysts, CERTs and agencies (NTRO context), critical-infrastructure operators — anyone with flow telemetry and alert fatigue"),
    (B, "•  Early, explained warnings redirect analyst attention to the right flows while an attack is still unfolding"),
    (H, "Benefits of the solution"),
    (B, "•  Economic: breaches average US$4.44M and ~241 days to identify and contain (IBM 2025); early warning + explanation shortens both; security AI/automation is associated with ~$1.9M lower breach costs"),
    (B, "•  Social: protects critical infrastructure and public services; metadata-only analysis preserves privacy (no packet payloads)"),
    (B, "•  Analyst welfare: explained, high-precision warnings instead of black-box alert floods"),
    (B, "•  Environmental: software-only, 57K-parameter model — no new hardware; edge deployment minimises data movement"),
    (B, "•  Scalable and globally comparable: trained/evaluated on the international benchmark dataset; the architecture (sensor → windows → model) is network-agnostic — any organization exporting flow data can deploy it"),
], top=Inches(1.30), height=Inches(5.7))
notes(S[4], (
    "IBM Cost of a Data Breach 2025 figures — verify at ibm.com/reports/"
    "data-breach before submitting. The 241-day identify+contain number is the "
    "product thesis in one stat: every day earlier is the entire value "
    "proposition."))

# ---------------------------------------------------- slide 6: references
tb = next(sh for sh in S[5].shapes if sh.name == "TextBox 8")
fill(tb, [
    (H, "Research and references"),
    (B, "•  CSE-CIC-IDS2018 dataset — Canadian Institute for Cybersecurity, UNB:  https://www.unb.ca/cic/datasets/ids-2018.html"),
    (B, "•  MITRE ATT&CK framework:  https://attack.mitre.org"),
    (B, "•  Sundararajan, Taly & Yan (2017), 'Axiomatic Attribution for Deep Networks' (Integrated Gradients), ICML"),
    (B, "•  Hochreiter & Schmidhuber (1997), 'Long Short-Term Memory', Neural Computation"),
    (B, "•  IBM Security, 'Cost of a Data Breach Report 2025':  https://www.ibm.com/reports/data-breach"),
    (B, "•  Open-source stack: Python, PyTorch, scapy, FastAPI, Next.js"),
    (B, "•  Project repository (working prototype + code):  <add GitHub link before submission>"),
], top=Inches(1.30), height=Inches(5.7))
notes(S[5], (
    "All benchmark numbers are reproducible from the repository "
    "(models/metrics_*.json). The demo video for the portal should follow "
    "docs/DEMO_RUNBOOK.md's 7-minute arc."))

# ------------------------------------------- slide 7: instructions — remove
sid = list(prs.slides._sldIdLst)[-1]
prs.part.drop_rel(sid.rId)
prs.slides._sldIdLst.remove(sid)

prs.save(str(DST))
print(f"wrote {DST} with {len(Presentation(str(DST)).slides)} slides")
