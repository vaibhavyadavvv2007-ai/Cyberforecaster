from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

p = Presentation('C:/Users/Abdul Fattah/Desktop/SIH2026-IDEA-Presentation-Format.pptx')
print(f'Slides: {len(p.slides)}')
print(f'Slide size: {p.slide_width} x {p.slide_height} ({p.slide_width/914400:.1f}\" x {p.slide_height/914400:.1f}\")')
print(f'Slide size in EMU: W={p.slide_width}, H={p.slide_height}')

for i, slide in enumerate(p.slides):
    print(f'\n=== Slide {i+1} ===')
    for j, shape in enumerate(slide.shapes):
        st = str(shape.shape_type)
        print(f'  [{j}] Shape: {st}, name="{shape.name}"')
        print(f'      pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        print(f'      left_in={shape.left/914400:.2f}, top_in={shape.top/914400:.2f}, w_in={shape.width/914400:.2f}, h_in={shape.height/914400:.2f}')
        try:
            if shape.has_text_frame:
                tf = shape.text_frame
                print(f'      Text frames: {len(tf.paragraphs)} paragraphs')
                for k, para in enumerate(tf.paragraphs):
                    text = para.text[:150] if para.text else '(empty)'
                    runs_info = []
                    for run in para.runs:
                        runs_info.append(f'font={run.font.name}, size={run.font.size}, bold={run.font.bold}, text="{run.text[:50]}"')
                    print(f'      Para {k}: [{text}]')
                    if runs_info:
                        for ri in runs_info:
                            print(f'        Run: {ri}')
        except Exception as e:
            print(f'      Error reading text: {e}')
        try:
            is_ph = shape.is_placeholder
            if is_ph:
                pf = shape.placeholder_format
                print(f'      IS PLACEHOLDER: type={pf.type}, idx={pf.idx}')
        except Exception:
            pass
